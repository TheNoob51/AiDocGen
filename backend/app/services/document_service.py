from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import os

def create_docx(project_data):
    document = Document()
    
    # Title
    document.add_heading(project_data.get("title", "Untitled Document"), 0)
    
    content = project_data.get("content", {})
    config = project_data.get("configuration", {})
    outline = config.get("outline", [])
    
    # If outline exists, follow it. Otherwise just dump content.
    sections_to_render = outline if outline else content.keys()
    
    for section in sections_to_render:
        document.add_heading(section, level=1)
        text = content.get(section, "")
        document.add_paragraph(text)
        
    # Save to buffer
    buffer = io.BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer

def create_pptx(project_data):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project_data.get("title", "Untitled Presentation")
    subtitle.text = project_data.get("description", "")
    
    content = project_data.get("content", {})
    config = project_data.get("configuration", {})
    slides_config = config.get("slides", [])
    
    # Iterate through configured slides
    for i, slide_config in enumerate(slides_config):
        slide_key = f"slide_{i}"
        slide_data = content.get(slide_key, {})
        
        # Bullet layout
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Title
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = slide_data.get("title", slide_config.get("title", f"Slide {i+1}"))
        
        # Content
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        
        bullets = slide_data.get("content", {}).get("bullets", [])
        if bullets:
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                
        # Notes
        notes = slide_data.get("content", {}).get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
            
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
